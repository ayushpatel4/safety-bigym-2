#!/usr/bin/env python
"""Build the FYP presentation deck (16:9 .pptx, editable in Keynote/PowerPoint).

Run:  ./venv/bin/python presentation/build_deck.py
Output: presentation/safety_bigym_presentation.pptx

Structure (interleaved slides + demo): title, motivation, gap+RQs, benchmark,
LIVE DEMO, architecture, two-axis protocol, RQ1 (persistent), RQ2 (intermittent),
RQ3 (regime map), conclusion, thanks  +  appendix slides for Q&A.

Speaker notes are embedded per slide; the full timed script lives in
presentation/speaker_script.md.
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ----------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent.parent          # safety_bigym/
FIG = ROOT / "FYP_final"
CLIP = ROOT / "presentation" / "assets" / "clips"
OUT = ROOT / "presentation" / "safety_bigym_presentation.pptx"

# palette
INK = RGBColor(0x1A, 0x1A, 0x1A)
NAVY = RGBColor(0x0F, 0x2A, 0x4A)
ORANGE = RGBColor(0xB4, 0x53, 0x09)     # persistent
BLUE = RGBColor(0x1E, 0x3A, 0x8A)       # intermittent
GREEN = RGBColor(0x14, 0x53, 0x2D)      # good / works
RED = RGBColor(0x8B, 0x1A, 0x1A)        # bad / fails
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Helvetica Neue"
MONO = "Menlo"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# helpers
def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    if bg is not None:
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        r.fill.solid(); r.fill.fore_color.rgb = bg
        r.line.fill.background()
        r.shadow.inherit = False
        # send to back
        sp = r._element
        sp.getparent().remove(sp)
        s.shapes._spTree.insert(2, sp)
    return s


def _set_font(run, size, bold, color, name=FONT, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name
    run.font.color.rgb = color


def textbox(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            wrap=True):
    """paras: list of dicts {text, size, bold, color, name, italic, bullet,
    level, space_before, space_after, align}. A para may also be a list of
    run-dicts for mixed styling on one line."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        runs = p if isinstance(p, list) else [p]
        meta = runs[0]
        para.alignment = meta.get("align", align)
        para.level = meta.get("level", 0)
        if meta.get("space_before") is not None:
            para.space_before = Pt(meta["space_before"])
        para.space_after = Pt(meta.get("space_after", 6))
        bullet = meta.get("bullet")
        prefix = ""
        if bullet:
            prefix = ("\u25B8 " if bullet is True else f"{bullet} ")
        for i, rd in enumerate(runs):
            r = para.add_run()
            r.text = (prefix if i == 0 else "") + rd["text"]
            _set_font(r, rd.get("size", 18), rd.get("bold", False),
                      rd.get("color", INK), rd.get("name", FONT),
                      rd.get("italic", False))
    return tb


def title_bar(s, title, kicker=None, color=NAVY):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set_font(r, 30, True, WHITE)
    # accent underline
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.12), SW, Pt(5))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background(); line.shadow.inherit = False
    if kicker:
        textbox(s, Inches(0.5), Inches(1.22), Inches(12.33), Inches(0.4),
                [{"text": kicker, "size": 15, "italic": True, "color": GREY}])
    return Inches(1.7 if kicker else 1.4)


_SLIDE_NO = 1  # title slide is page 1 (carries no footer)


def footer(s, idx=None, short=""):
    """idx is ignored for main slides (auto-numbered); pass "A" for appendix."""
    global _SLIDE_NO
    if idx == "A":
        label = "A"
    else:
        _SLIDE_NO += 1
        label = str(_SLIDE_NO)
    textbox(s, Inches(0.5), Inches(7.06), Inches(9), Inches(0.35),
            [{"text": f"safety_bigym  \u00b7  Ayush Patel  \u00b7  {short}",
              "size": 11, "color": GREY}])
    textbox(s, Inches(12.0), Inches(7.06), Inches(0.9), Inches(0.35),
            [{"text": label, "size": 11, "color": GREY, "align": PP_ALIGN.RIGHT}])


def img_size(path):
    with Image.open(path) as im:
        return im.size


def add_image_fit(s, path, x, y, max_w, max_h, align="center", valign="top"):
    iw, ih = img_size(path)
    ar = iw / ih
    w = max_w; h = Emu(int(w / ar))
    if h > max_h:
        h = max_h; w = Emu(int(h * ar))
    if align == "center":
        x = Emu(int(x + (max_w - w) / 2))
    elif align == "right":
        x = Emu(int(x + (max_w - w)))
    if valign == "middle":
        y = Emu(int(y + (max_h - h) / 2))
    pic = s.shapes.add_picture(str(path), x, y, width=w, height=h)
    pic.line.color.rgb = RGBColor(0xD0, 0xD7, 0xDE)
    pic.line.width = Pt(0.75)
    return pic, w, h


def add_movie(s, mp4, x, y, max_w, max_h):
    poster = Path(mp4).with_suffix(".png")
    iw, ih = img_size(poster)
    ar = iw / ih
    w = max_w; h = Emu(int(w / ar))
    if h > max_h:
        h = max_h; w = Emu(int(h * ar))
    x = Emu(int(x + (max_w - w) / 2))
    mv = s.shapes.add_movie(str(mp4), x, y, w, h,
                            poster_frame_image=str(poster), mime_type="video/mp4")
    return mv, w, h


def panel(s, x, y, w, h, fill=LIGHT, line=None, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    return shp


def chip(s, x, y, w, text, fill, tcolor=WHITE, h=Inches(0.42), size=14):
    c = panel(s, x, y, w, h, fill=fill)
    tf = c.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.08)
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    _set_font(r, size, True, tcolor)
    return c


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()


def big_stat(s, x, y, w, value, label, vcolor=ACCENT):
    panel(s, x, y, w, Inches(1.5), fill=LIGHT)
    textbox(s, x, y + Inches(0.16), w, Inches(0.8),
            [{"text": value, "size": 34, "bold": True, "color": vcolor,
              "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    textbox(s, x, y + Inches(0.98), w, Inches(0.45),
            [{"text": label, "size": 13, "color": INK, "align": PP_ALIGN.CENTER}],
            align=PP_ALIGN.CENTER)


def arrow(s, x1, y1, x2, y2, color=NAVY, weight=2.25):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = Pt(weight)
    try:
        cn.line.headEnd = None
    except Exception:
        pass
    return cn


# ============================================================================
# SLIDE 1 — Title
# ============================================================================
s = slide(NAVY)
# subtle accent block
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.55), SW, Pt(6))
band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()
band.shadow.inherit = False
textbox(s, Inches(0.9), Inches(1.45), Inches(11.6), Inches(3.0), [
    {"text": "safety_bigym", "size": 30, "bold": True, "color": RGBColor(0x9E,0xC5,0xFF), "space_after": 10},
    [{"text": "A Co-location Regime Map for", "size": 39, "bold": True, "color": WHITE, "space_after": 2}],
    [{"text": "Safe Humanoid Manipulation", "size": 39, "bold": True, "color": WHITE, "space_after": 2}],
    [{"text": "Alongside a Live Human Coworker", "size": 39, "bold": True, "color": WHITE}],
])
textbox(s, Inches(0.9), Inches(5.75), Inches(11.5), Inches(1.3), [
    {"text": "Ayush Patel", "size": 22, "bold": True, "color": WHITE},
    {"text": "MEng Final Year Project  \u00b7  Supervisor: Stephen James", "size": 17, "color": RGBColor(0xC9,0xD6,0xE5), "space_before": 2},
    {"text": "Department of Computing, Imperial College London", "size": 17, "color": RGBColor(0xC9,0xD6,0xE5)},
])
notes(s, """
[0:00-0:15] Good morning. I'm Ayush. My project asks a safety question that is
becoming urgent: when a *learned* humanoid robot has to work in the same space a
person is reaching into, what actually keeps that person safe? Over the next ~16
minutes I'll show you the benchmark I built, the safety architecture, and the
one finding I want you to remember - which I'll demo live.
""")

# ============================================================================
# SLIDE 2 — Why this matters
# ============================================================================
s = slide()
top = title_bar(s, "Robot learning has left the cage")
textbox(s, Inches(0.5), top, Inches(6.7), Inches(5.0), [
    {"text": "Learned humanoids now share the workspace by design.", "size": 21, "bold": True, "color": NAVY, "space_after": 10},
    {"text": "Imitation learning and demo-driven RL put commercial humanoids (Unitree G1 / H1) on tasks performed right next to people.", "size": 17, "bullet": True, "space_after": 10},
    {"text": "But their policies optimise task success only.", "size": 18, "bold": True, "color": RED, "bullet": True, "space_after": 4},
    {"text": "They do not see an ISO 15066 safety margin, never reason about contact force, and often do not observe the human at all.", "size": 16, "level": 1, "space_after": 12, "color": INK},
    {"text": "Safety validation has not kept up.", "size": 18, "bold": True, "color": NAVY, "bullet": True, "space_after": 4},
    {"text": "No commercial humanoid was validated for collaborative operation by 2026; the H1 still needs a cage for lively testing.", "size": 16, "level": 1, "color": INK},
])
panel(s, Inches(7.5), top, Inches(5.3), Inches(4.9), fill=LIGHT)
add_image_fit(s, CLIP / "clip_benchmark_disruption.png", Inches(7.7), top + Inches(0.2),
              Inches(4.9), Inches(2.5))
textbox(s, Inches(7.7), top + Inches(2.85), Inches(4.9), Inches(2.0), [
    {"text": "ISO 15066 (now in ISO 10218:2025) defines two mechanisms for working near people:", "size": 14, "color": INK, "space_after": 6},
    {"text": "Speed & Separation Monitoring \u2013 keep a velocity-dependent distance.", "size": 14, "bullet": True, "color": NAVY, "space_after": 4},
    {"text": "Power & Force Limiting \u2013 cap contact force per body region.", "size": 14, "bullet": True, "color": NAVY},
])
footer(s, 2, "Motivation")
notes(s, """
[0:15-1:15] Two facts set up the whole project. First, robot learning has left
the cage - commercial humanoids are aimed at tasks performed right next to a
person. Second, the policies that drive them are trained to maximise task
reward; they do not optimise human distance, an ISO safety margin, or contact
force, and often don't even observe the human. There is a mature safety standard
for exactly this - ISO 15066, speed-separation and force-limiting - but learned
policies target none of it, and no commercial humanoid is validated for
collaborative operation. That gap is my starting point.
""")

# ============================================================================
# SLIDE 3 — The gap + research questions
# ============================================================================
s = slide()
top = title_bar(s, "The gap, and three questions",
                "Manipulation benchmarks ignore safety; safe-RL ignores humanoid manipulation beside a person")
# gap panels
gx, gw, gy = Inches(0.5), Inches(3.9), top + Inches(0.1)
for i, (t, sub, col) in enumerate([
    ("Humanoid manipulation", "BiGym, HumanoidBench \u2014 no safety, no human in the scene", BLUE),
    ("ISO-style safety / safe-RL", "Safety-Gymnasium \u2014 abstract costs, low-DoF, ground-truth state", ORANGE),
    ("Imperfect human perception", "real robots see noisy on-robot pose estimates, not mocap", GREEN),
]):
    px = Inches(0.5 + i * 4.27)
    panel(s, px, gy, gw, Inches(1.7), fill=LIGHT, line=col)
    textbox(s, px + Inches(0.18), gy + Inches(0.16), gw - Inches(0.36), Inches(1.4), [
        {"text": t, "size": 16, "bold": True, "color": col, "space_after": 5},
        {"text": sub, "size": 13.5, "color": INK},
    ])
textbox(s, Inches(0.5), gy + Inches(1.85), Inches(12.3), Inches(0.5), [
    {"text": "No prior work studies all three together on a high-DoF humanoid manipulator. This project does.",
     "size": 17, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
# RQ panel
ry = gy + Inches(2.5)
panel(s, Inches(0.5), ry, Inches(12.33), Inches(2.0), fill=NAVY)
textbox(s, Inches(0.8), ry + Inches(0.16), Inches(11.8), Inches(1.8), [
    [{"text": "RQ1  ", "size": 17, "bold": True, "color": RGBColor(0x9E,0xC5,0xFF)},
     {"text": "Training-time: can constrained RL reduce how often robot and human are too close, while still doing the task?", "size": 16, "color": WHITE}],
    [{"text": "RQ2  ", "size": 17, "bold": True, "color": RGBColor(0x9E,0xC5,0xFF), "space_before": 8},
     {"text": "Runtime: can a filter cut the robot's unsafe speed near the human without breaking the task?", "size": 16, "color": WHITE}],
    [{"text": "RQ3  ", "size": 17, "bold": True, "color": RGBColor(0x9E,0xC5,0xFF), "space_before": 8},
     {"text": "Choosing: can the human\u2013robot co-location pattern predict which of the two to use?", "size": 16, "color": WHITE}],
])
footer(s, 3, "The gap & RQs")
notes(s, """
[1:15-2:15] The gap is an intersection. Manipulation benchmarks like BiGym have
no safety and usually no human in the scene. Safe-RL benchmarks use abstract
costs, low-DoF systems, and assume perfect state. And almost none model the
imperfect human perception a real robot actually has. Nobody studies all three
together on a high-DoF humanoid. So I ask three questions: RQ1, can constrained
RL keep the robot away while still doing the task; RQ2, can a runtime filter cut
unsafe speed near the human; RQ3 - the interesting one - can a measurable
property of the task tell you which of those two you should reach for. Keep these
three in mind; the talk answers them in order.
""")

# ============================================================================
# SLIDE 4 — The benchmark
# ============================================================================
s = slide()
top = title_bar(s, "Contribution 1: the safety_bigym benchmark",
                "A safety-aware extension of BiGym for manipulation beside a moving coworker")
add_image_fit(s, FIG / "task_suite.png", Inches(0.5), top, Inches(12.33), Inches(3.0))
yb = top + Inches(3.15)
for i, (t, sub, col) in enumerate([
    ("Moving G1 coworker", "a physically-grounded human stand-in; benchmark is body-agnostic", NAVY),
    ("ISO 15066 metrics", "per-joint separation + speed-scaling, computed every step", NAVY),
    ("Calibrated perception noise", "mock-BodySLAM++: the robot sees a noisy human estimate", NAVY),
    ("Two co-location regimes", "persistent (saucepan) vs intermittent (dishwasher, drawers)", ACCENT),
]):
    px = Inches(0.5 + i * 3.12)
    panel(s, px, yb, Inches(2.95), Inches(1.65), fill=LIGHT, line=(ACCENT if i==3 else RGBColor(0xD0,0xD7,0xDE)))
    textbox(s, px + Inches(0.15), yb + Inches(0.14), Inches(2.65), Inches(1.4), [
        {"text": t, "size": 14.5, "bold": True, "color": col, "space_after": 4},
        {"text": sub, "size": 12.5, "color": INK},
    ])
footer(s, 4, "Benchmark")
notes(s, """
[2:15-3:30] My first contribution is the benchmark itself. safety_bigym extends
BiGym with: a moving Unitree-G1 coworker as a human stand-in - it's physically
grounded and the benchmark is body-agnostic, so this is about co-location
dynamics, not modelling a specific person; ISO-15066-derived metrics computed
every step on the closest human-joint / robot-link pair; calibrated perception
noise, so the policy sees a *noisy* human estimate like a real on-robot pose
estimator would produce; and crucially a three-task suite spanning two regimes -
saucepan_to_hob, where the human is in the workspace almost continuously, versus
dishwasher and drawers, where encounters come in bursts. That regime distinction
is the spine of the whole talk. Rather than just describe it, let me show it.
""")

# ============================================================================
# SLIDE 5 — LIVE DEMO marker
# ============================================================================
s = slide(NAVY)
chip(s, Inches(0.6), Inches(0.55), Inches(3.0), "LIVE DEMONSTRATION", ACCENT, WHITE, h=Inches(0.55), size=18)
textbox(s, Inches(0.6), Inches(1.4), Inches(6.4), Inches(5.2), [
    {"text": "The benchmark, running live", "size": 30, "bold": True, "color": WHITE, "space_after": 14},
    {"text": "G1 coworker parks in the robot's working volume and reaches at it on a cycle \u2014 min separation ~0.15 m.", "size": 17, "bullet": True, "color": WHITE, "space_after": 8},
    {"text": "Terminal prints each reach phase (extend / hold / retract), the target (robot end-effector), and the reach gate, live.", "size": 17, "bullet": True, "color": WHITE, "space_after": 8},
    {"text": "\u201cThis is the persistent regime \u2013 the human is in the way most of the time.\u201d", "size": 16, "italic": True, "color": RGBColor(0x9E,0xC5,0xFF), "space_after": 14},
])
panel(s, Inches(0.6), Inches(5.55), Inches(6.4), Inches(1.25), fill=RGBColor(0x0A,0x1F,0x38))
textbox(s, Inches(0.8), Inches(5.66), Inches(6.0), Inches(1.1), [
    {"text": "cd safety_bigym", "size": 13, "color": RGBColor(0x9E,0xC5,0xFF), "name": MONO, "space_after": 2},
    {"text": "export AMASS_DATA_DIR=.../CMU/CMU", "size": 13, "color": RGBColor(0x9E,0xC5,0xFF), "name": MONO, "space_after": 2},
    {"text": "./venv/bin/mjpython scripts/demo_coworker.py --human g1 \\", "size": 12.5, "color": WHITE, "name": MONO, "space_after": 2},
    {"text": "   --spawn in_place --reach-target ee --stage train --task saucepan", "size": 12.5, "color": WHITE, "name": MONO},
])
# backup video on the right
textbox(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(0.4),
        [{"text": "Backup recording (click to play if live fails):", "size": 13, "italic": True, "color": RGBColor(0xC9,0xD6,0xE5)}])
add_movie(s, CLIP / "clip_benchmark_disruption.mp4", Inches(7.2), Inches(1.85),
          Inches(5.6), Inches(4.8))
footer(s, 5, "Live demo")
notes(s, """
[3:30-6:30] LIVE DEMO (target ~3 min). Switch to the MuJoCo window already
running. Narrate: the tan figure is the coworker stand-in; the dark robot is the
H1 manipulator. Watch the coworker walk in and start reaching toward the robot's
end-effector. Flip to the terminal: it prints the closest human-joint to
robot-link pair and the live separation - point out the closest joint changing
(wrist, elbow) as the arm sweeps. Key line: 'this is the persistent regime - the
human is in the way most of the time, which is going to matter enormously in a
minute.' If the viewer misbehaves, click the backup recording on this slide.
Then: 'training the policies takes ~450 GPU-hours on A100s and the checkpoints
live on the lab GPU box, so the trained behaviours that follow are recorded
rollouts - exactly how my evaluation harness produces them.'
""")

# ============================================================================
# SLIDE 6 — Why these tasks + curriculum
# ============================================================================
s = slide()
top = title_bar(s, "Why these three tasks \u2014 and why a curriculum",
                "Task choice IS the experiment; curriculum is what makes the task learnable")
# left: task selection
panel(s, Inches(0.5), top, Inches(6.05), Inches(3.15), fill=LIGHT, line=ACCENT)
textbox(s, Inches(0.7), top + Inches(0.14), Inches(5.7), Inches(2.9), [
    {"text": "Tasks span the regime axis", "size": 17, "bold": True, "color": ACCENT, "space_after": 5},
    {"text": "the thesis's independent variable, not an afterthought", "size": 12.5, "italic": True, "color": GREY, "space_after": 8},
    [{"text": "Persistent: ", "size": 14.5, "bold": True, "color": ORANGE},
     {"text": "saucepan_to_hob \u2014 human + robot share the hob/counter almost continuously.", "size": 14, "color": INK}],
    {"text": "", "size": 3},
    [{"text": "Intermittent: ", "size": 14.5, "bold": True, "color": BLUE},
     {"text": "dishwasher_close, drawers_open_all \u2014 robot works at an appliance; the human passes through in bursts.", "size": 14, "color": INK}],
    {"text": "", "size": 3},
    {"text": "All are BiGym kitchen tasks with expert demos \u2014 required by the demo-driven CQN-AS backbone.", "size": 13.5, "bullet": True, "color": INK},
])
# right: curriculum
panel(s, Inches(6.78), top, Inches(6.05), Inches(3.15), fill=LIGHT, line=ORANGE)
textbox(s, Inches(6.98), top + Inches(0.14), Inches(5.7), Inches(2.9), [
    {"text": "Why a curriculum", "size": 17, "bold": True, "color": ORANGE, "space_after": 5},
    [{"text": "Direct training on the full coworker fails: ", "size": 14.5, "bold": True, "color": RED},
     {"text": "avoidance dominates the reward and the manipulation task is never discovered.", "size": 14, "color": INK}],
    {"text": "", "size": 3},
    {"text": "Fix: stage the coworker's behaviour from easy to hard, so task competence is bootstrapped first.", "size": 14, "bullet": True, "color": INK, "space_after": 4},
    {"text": "Standard for long-horizon sparse-reward RL \u2014 and reported as a limitation (results presuppose it).", "size": 13.5, "bullet": True, "italic": True, "color": INK},
])
# bottom: 3-stage curriculum chips
cy = top + Inches(3.45)
stages = [("Stage 0 \u00b7 idle", "coworker distant & still", RGBColor(0x64,0x74,0x8B)),
          ("Stage 1 \u00b7 easy", "gentle, occasional reaches", ACCENT),
          ("Stage 2 \u00b7 full", "coworker_train: close, frequent reaches", ORANGE)]
cw = Inches(3.85)
xs3 = [Inches(0.5), Inches(4.74), Inches(8.98)]
for i, (t, sub, col) in enumerate(stages):
    panel(s, xs3[i], cy, cw, Inches(1.15), fill=col)
    textbox(s, xs3[i] + Inches(0.15), cy + Inches(0.16), cw - Inches(0.3), Inches(0.95), [
        {"text": t, "size": 15, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "space_after": 3},
        {"text": sub, "size": 12, "color": WHITE, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    if i < 2:
        ax = int(xs3[i]) + int(cw)
        arrow(s, Emu(ax), cy + Inches(0.57), xs3[i+1], cy + Inches(0.57), color=INK, weight=3)
footer(s, short="Tasks & curriculum")
notes(s, """
[after the demo] Two questions before the methods. First, why these three tasks?
The task choice IS the experiment: I picked tasks that span the co-location axis -
saucepan_to_hob, where the human and robot share the hob and counter almost
continuously (persistent), versus dishwasher and drawers, where the robot works
at an appliance and the human only passes through in bursts (intermittent). They're
all BiGym kitchen tasks with expert demonstrations, which the demo-driven CQN-AS
backbone needs. Second, why a curriculum? If you train directly on the full
coworker disruption, avoidance dominates the reward and the robot never discovers
the manipulation task at all. So I stage the coworker's behaviour - idle, then
easy, then the full coworker_train band - which bootstraps task competence first.
It's a standard fix for long-horizon sparse-reward RL, and I report the dependence
on it as a limitation.
""")

# ============================================================================
# SLIDE 7 — Architecture (two arms)
# ============================================================================
s = slide()
top = title_bar(s, "Contribution 2: the Hybrid Safety Critic",
                "Two mechanisms with different jobs \u2014 a proactive policy and a reactive runtime filter")

def pipe_box(s, x, y, w, h, title, sub, fill, tcol=WHITE, subcol=None):
    b = panel(s, x, y, w, h, fill=fill, radius=True)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; _set_font(r, 14.5, True, tcol)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        _set_font(r2, 11.5, False, subcol or tcol, name=MONO)
    return b

row_y = top + Inches(0.4)
bh = Inches(1.15)
xs = [Inches(0.5), Inches(2.7), Inches(5.1), Inches(8.0), Inches(11.0)]
ws = [Inches(1.9), Inches(2.1), Inches(2.6), Inches(2.6), Inches(1.8)]
pipe_box(s, xs[0], row_y, ws[0], bh, "Environment", "BiGym + G1", RGBColor(0x37,0x47,0x55))
pipe_box(s, xs[1], row_y, ws[1], bh, "Observation", "proprio + noisy p\u0302", RGBColor(0x37,0x47,0x55))
pipe_box(s, xs[2], row_y, ws[2], bh, "Policy  \u03c0", "argmax  Qr \u2212 \u03bb\u00b7Qc", BLUE, WHITE, RGBColor(0xCF,0xE0,0xFF))
pipe_box(s, xs[3], row_y, ws[3], bh, "Runtime filter", "Qsafe(s,a) \u2265 R ?", ORANGE, WHITE, RGBColor(0xFF,0xE6,0xC7))
pipe_box(s, xs[4], row_y, ws[4], bh, "Execute  a", "", GREEN)
midy = int(row_y) + int(bh) // 2
for i in range(4):
    arrow(s, xs[i] + ws[i], Emu(midy), xs[i+1], Emu(midy), color=NAVY)
# training bands
tb_y = row_y + Inches(1.7)
panel(s, xs[2], tb_y, ws[2], Inches(0.95), fill=LIGHT, line=BLUE)
textbox(s, xs[2] + Inches(0.1), tb_y + Inches(0.1), ws[2] - Inches(0.2), Inches(0.8), [
    {"text": "trained ONLINE", "size": 12, "bold": True, "color": BLUE, "align": PP_ALIGN.CENTER, "space_after": 2},
    {"text": "value-based Lagrangian on cost c\u209c \u2192 Qr, Qc", "size": 11.5, "color": INK, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
panel(s, xs[3], tb_y, ws[3], Inches(0.95), fill=LIGHT, line=ORANGE)
textbox(s, xs[3] + Inches(0.1), tb_y + Inches(0.1), ws[3] - Inches(0.2), Inches(0.8), [
    {"text": "trained OFFLINE (CQL), frozen", "size": 12, "bold": True, "color": ORANGE, "align": PP_ALIGN.CENTER, "space_after": 2},
    {"text": "safety value function \u2192 Qsafe", "size": 11.5, "color": INK, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
arrow(s, xs[2] + Emu(int(ws[2]/2)), tb_y, xs[2] + Emu(int(ws[2]/2)), row_y + bh, color=BLUE, weight=1.75)
arrow(s, xs[3] + Emu(int(ws[3]/2)), tb_y, xs[3] + Emu(int(ws[3]/2)), row_y + bh, color=ORANGE, weight=1.75)
# two-job summary
jy = tb_y + Inches(1.25)
panel(s, Inches(0.5), jy, Inches(6.05), Inches(1.45), fill=LIGHT, line=BLUE)
textbox(s, Inches(0.7), jy + Inches(0.14), Inches(5.7), Inches(1.2), [
    {"text": "Policy = PROACTIVE", "size": 15, "bold": True, "color": BLUE, "space_after": 3},
    {"text": "learns to keep distance during training (internalised, anticipatory).", "size": 13.5, "color": INK},
])
panel(s, Inches(6.78), jy, Inches(6.05), Inches(1.45), fill=LIGHT, line=ORANGE)
textbox(s, Inches(6.98), jy + Inches(0.14), Inches(5.7), Inches(1.2), [
    {"text": "Filter = REACTIVE", "size": 15, "bold": True, "color": ORANGE, "space_after": 3},
    {"text": "a deployment-time check: veto, dodge, speed-scale, or critic-gate.", "size": 13.5, "color": INK},
])
footer(s, 6, "Architecture")
notes(s, """
[6:30-7:45] The architecture has two arms with genuinely different jobs. At each
step the policy proposes an action by maximising task value minus lambda times
cost value - that's the constrained, *proactive* part, trained online; it learns
to keep its distance. Then a runtime *filter* checks that action against a
learned safety value function and can pass it, slow it, dodge, or veto - that's
the *reactive* part, trained offline and frozen. One implementation note for the
specialists: my backbone, CQN-AS, is critic-only - there's no actor network - so
the usual actor-critic Lagrangian had to be re-derived in value-based form as a
dual-Q argmax, and I proved a C51 value-support bound so dense safety shaping
doesn't silently saturate the critic. That re-derivation is one of my
contributions. The headline question is: which arm should carry safety?
""")

# ============================================================================
# SLIDE 7 — Two-axis protocol
# ============================================================================
s = slide()
top = title_bar(s, "Contribution 3: measuring safety honestly",
                "Two axes \u2014 and most closeness is the human's doing, not the robot's")
add_image_fit(s, FIG / "fig4_exogenous_proximity.png", Inches(6.7), top, Inches(6.1), Inches(4.7))
textbox(s, Inches(0.5), top + Inches(0.1), Inches(6.0), Inches(5.0), [
    {"text": "Two safety axes:", "size": 19, "bold": True, "color": NAVY, "space_after": 8},
    [{"text": "Proximity ", "size": 16.5, "bold": True, "color": INK},
     {"text": "\u2013 how often they are too close (exposure).", "size": 16.5, "color": INK, "bullet": False}],
    {"text": "", "size": 4},
    [{"text": "Velocity-adaptive SSM ", "size": 16.5, "bold": True, "color": INK},
     {"text": "\u2013 is the robot moving too fast for the current separation (robot-controllable).", "size": 16.5, "color": INK}],
    {"text": "Freeze the robot completely \u2192 proximity drops only ~58%.", "size": 17, "bold": True, "color": RED, "bullet": True, "space_before": 14, "space_after": 4},
    {"text": "~42% of closeness is the human walking into a stationary robot \u2014 exogenous, uncontrollable.", "size": 15.5, "level": 1, "color": INK, "space_after": 12},
    {"text": "So the robot's own safety is judged mostly on the velocity axis \u2014 and the worst-case tail is reported, not hidden.", "size": 16, "italic": True, "color": NAVY, "bullet": True},
])
footer(s, 7, "Two-axis protocol")
notes(s, """
[7:45-8:45] Before any results, a measurement point that earns its slide. I split
safety into two axes: proximity - how often they're too close, the exposure - and
velocity-adaptive SSM - whether the robot is moving too fast for the current gap,
which is the robot-controllable part. Why split? Because I ran a frozen-robot
sweep: even with the robot completely stopped, proximity only falls about 58%.
Roughly 42% of all closeness is just the human walking into a stationary robot -
it's exogenous, no robot policy can remove it. That single fact reframes
everything: I judge the robot's safety mainly on the velocity axis, and I report
the uncontrollable worst-case tail honestly rather than claiming credit I can't
earn. Now the results, in three movements.
""")

# ============================================================================
# SLIDE 8 — RQ1 persistent
# ============================================================================
s = slide()
top = title_bar(s, "RQ1 \u00b7 Persistent co-location: train safety in", color=ORANGE)
# video
textbox(s, Inches(0.5), top, Inches(7.4), Inches(0.35),
        [{"text": "saucepan_to_hob  \u00b7  baseline vs constrained  (scripted reconstruction \u2014 real env / coworker / ISO monitor; HUD measured)", "size": 11.5, "italic": True, "color": GREY}])
add_movie(s, CLIP / "clip_avoid_compare.mp4", Inches(0.5), top + Inches(0.42),
          Inches(7.5), Inches(3.45))
# stats
sx = Inches(8.2)
big_stat(s, sx, top + Inches(0.1), Inches(4.55), "0.296 \u2192 0.228", "proximity-violation rate  (\u221223%, 3 seeds)", GREEN)
big_stat(s, sx, top + Inches(1.75), Inches(4.55), "0.76", "task success retained (baseline 0.85)", NAVY)
textbox(s, Inches(0.5), top + Inches(3.55), Inches(12.3), Inches(2.0), [
    {"text": "The constrained policy is the only mechanism that reduces exposure gracefully \u2014 it yields as the coworker closes, then returns.", "size": 16, "bold": True, "color": NAVY, "space_after": 8},
    [{"text": "Both axes at once ", "size": 15.5, "bold": True, "color": INK},
     {"text": "needs policy + speed-scaling \u2014 the regime's ceiling, at a real cost (success 0.85 \u2192 0.44). Reported, not hidden.", "size": 15.5, "color": INK}],
    [{"text": "Method gem: ", "size": 15.5, "bold": True, "color": ACCENT},
     {"text": "picking the checkpoint by peak success picks AGAINST the constraint \u2014 the avoidance lives in a mid-training basin (two false nulls before I found it).", "size": 15.5, "color": INK}],
])
footer(s, short="RQ1 \u00b7 persistent")
notes(s, """
[~9:00-10:00] RQ1, the persistent task. Play the side-by-side. Left, red: the
unconstrained baseline works straight through the coworker's reach window -
repeated violations, min sep ~0.2 m. Right, green: the constrained policy, fixed
lambda = 0.1, yields away and returns once the coworker leaves. The numbers:
proximity-violation rate 0.296 down to 0.228, a 23% cut, reproducible across
three seeds, at 0.76 success against a 0.85 baseline. Getting *both* safety axes
low needs the policy composed with speed-scaling, and that costs success, 0.85 to
0.44 - I report that ceiling rather than hide it. One honesty gem the markers
like: selecting the checkpoint by peak success picks *against* the safety
constraint; the avoidance lives in a mid-training basin, and selecting on success
gave me two false nulls before I found it. (Note for honesty: these clips are
scripted reconstructions - real environment, coworker and ISO monitor, scripted
robot - because the trained snapshots live on the GPU box; every HUD number is
measured from the sim.) Transition: 'a fair question is - why not just bolt a
safety filter onto the baseline?'
""")

# ============================================================================
# SLIDE 10 — RQ1 reactive filters freeze/flee (filter video)
# ============================================================================
s = slide()
top = title_bar(s, "RQ1 \u00b7 Why not just bolt on a filter? Freeze vs flee",
                "A reactive filter acts only once the human is already close \u2014 then it can only freeze or flee", color=ORANGE)
textbox(s, Inches(0.5), top, Inches(8.0), Inches(0.32),
        [{"text": "saucepan_to_hob  \u00b7  learned-veto fallbacks  (scripted reconstruction; HUD measured from sim)", "size": 11.5, "italic": True, "color": GREY}])
add_movie(s, CLIP / "clip_veto_compare.mp4", Inches(0.5), top + Inches(0.4),
          Inches(7.7), Inches(3.5))
panel(s, Inches(8.45), top + Inches(0.1), Inches(4.4), Inches(3.7), fill=LIGHT)
textbox(s, Inches(8.65), top + Inches(0.28), Inches(4.05), Inches(3.4), [
    {"text": "FREEZE (zero-velocity veto)", "size": 15, "bold": True, "color": RED, "space_after": 3},
    {"text": "robot stops and dwells in the danger zone the coworker creates \u2014 proximity unchanged (0.296 \u2192 0.303).", "size": 13.5, "color": INK, "space_after": 10},
    {"text": "FLEE (retreat veto)", "size": 15, "bold": True, "color": ORANGE, "space_after": 3},
    {"text": "buys separation (0.296 \u2192 0.095) only by abandoning the task: success 0.85 \u2192 0.18, mean velocity \u00d76.", "size": 13.5, "color": INK, "space_after": 10},
    {"text": "Same limit reached by the CBF / safety-filter literature. Anticipation must come from the policy.", "size": 13.5, "italic": True, "color": NAVY},
])
textbox(s, Inches(0.5), top + Inches(4.05), Inches(12.3), Inches(0.8), [
    {"text": "No reactive filter gracefully reduces exposure under persistent co-location \u2014 which is exactly why RQ1's answer is to train safety into the policy.",
     "size": 16, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
footer(s, short="RQ1 \u00b7 freeze vs flee")
notes(s, """
[~10:00-11:00] The obvious objection: why not just bolt a safety filter onto the
baseline? This is the answer, and it's half the thesis. A reactive filter only
acts once the human is *already* close - at which point the robot has two options.
Left: FREEZE - the zero-velocity veto stops the robot, but it then just dwells in
the danger zone the coworker is walking into, so proximity is unchanged, 0.296 to
0.303. Right: FLEE - the retreat veto buys distance, 0.296 down to 0.095, but only
by abandoning the task: success collapses 0.85 to 0.18 and mean velocity goes up
six-fold, which is itself an SSM hazard. Every reactive modality I tried, learned
or model-based, hits this freeze-versus-flee wall - a limit the control-barrier
literature reaches independently. The missing ingredient is *anticipation*, and
that has to be trained into the policy. So under persistent co-location, safety
must be trained in. (Scripted reconstruction; HUD measured from sim.)
""")

# ============================================================================
# SLIDE 11 — RQ2 intermittent
# ============================================================================
s = slide()
top = title_bar(s, "RQ2 \u00b7 Intermittent co-location: gate the backstop", color=BLUE)
add_image_fit(s, FIG / "fig1_method_comparison.png", Inches(0.5), top, Inches(7.0), Inches(3.45))
textbox(s, Inches(7.55), top, Inches(5.3), Inches(0.3),
        [{"text": "ISO-SSM speed-scaling  (scripted reconstruction; HUD measured)", "size": 11, "italic": True, "color": GREY}])
add_movie(s, CLIP / "clip_speedscale_compare.mp4", Inches(7.55), top + Inches(0.34),
          Inches(5.3), Inches(3.15))
textbox(s, Inches(0.5), top + Inches(3.7), Inches(12.3), Inches(1.85), [
    [{"text": "Constrained RL finds no feasible budget here ", "size": 15.5, "bold": True, "color": RED},
     {"text": "\u2014 and when \u03bb binds, the robot gets FASTER (0.44 \u2192 0.79 m/s), not safer. A binary veto breaks the chunked policy (max vel 2.46 \u2192 6.03 m/s).", "size": 15.5, "color": INK}],
    [{"text": "Winner: a learned critic GATING a graded ISO speed-scaler. ", "size": 15.5, "bold": True, "color": GREEN},
     {"text": "Critic decides WHEN (fire if Qsafe < R); scaler decides HOW (slow smoothly).", "size": 15.5, "color": INK}],
    [{"text": "dishwasher \u221250% SSM at \u22120.10 success  \u00b7  drawers \u221222% at \u22120.09 (or \u221210% at \u22120.02).  ", "size": 15.5, "bold": True, "color": NAVY},
     {"text": "Gate threshold R is a deployable dial.", "size": 15.5, "italic": True, "color": INK}],
])
footer(s, short="RQ2 \u00b7 intermittent")
notes(s, """
[10:15-11:45] RQ2, the intermittent tasks - dishwasher and drawers. Here the
story flips. Constrained RL finds *no feasible budget*: the constraint is inert
above the natural cost and task-fatal below it, and when the multiplier does bind
the robot actually gets faster, not safer. A from-scratch WCSAC baseline
corroborates that trainability difficulty. And a binary veto shatters the chunked
policy - max velocity jumps from 2.5 to 6 metres per second. What works is a
learned critic *gating* a graded ISO speed-scaler: the critic decides WHEN to
intervene, the scaler decides HOW, slowing smoothly. On dishwasher that's a 50%
cut in SSM violations for a 10-point success cost; on drawers, 22% - or a nearly
free 10% cut at almost no cost. The clip on the right shows the law in action:
full speed while the coworker is clear, graded slow-down as separation closes,
SSM-OK throughout. And the gate threshold R is a clean dial you tune at deployment.
""")

# ============================================================================
# SLIDE 10 — RQ3 the regime map (payoff)
# ============================================================================
s = slide()
top = title_bar(s, "RQ3 \u00b7 The regime map \u2014 and a test fixed in advance",
                "The deciding variable is measurable: the fraction of steps the safety gate is active")
add_image_fit(s, FIG / "fig7_gate_activity.png", Inches(0.5), top, Inches(6.0), Inches(3.5))
add_image_fit(s, FIG / "fig5_cross_task_boundary.png", Inches(6.8), top, Inches(6.0), Inches(3.5))
panel(s, Inches(0.5), top + Inches(3.7), Inches(12.33), Inches(1.55), fill=LIGHT, line=ACCENT)
textbox(s, Inches(0.75), top + Inches(3.83), Inches(11.9), Inches(1.4), [
    [{"text": "Gate-active fraction IS the regime:  ", "size": 16, "bold": True, "color": NAVY},
     {"text": "saucepan 61.5% of steps  vs  dishwasher 26.5% / drawers 19.2%  (matched R).", "size": 16, "color": INK}],
    [{"text": "Pre-registered rule (fixed before running): ", "size": 15.5, "bold": True, "color": ACCENT},
     {"text": "gating \u201crecovers throughput\u201d iff some R reaches success \u2265 0.60 AND SSM \u2264 0.08.", "size": 15.5, "color": INK}],
    [{"text": "On saucepan, NO row passed \u2014 exactly the predicted failure. ", "size": 15.5, "bold": True, "color": RED},
     {"text": "Intermittent tasks bend into the safe corner; the persistent task slides down the diagonal.", "size": 15.5, "italic": True, "color": INK}],
])
footer(s, 10, "RQ3 \u00b7 regime map")
notes(s, """
[11:45-13:15] RQ3 ties the two halves into one finding and validates it. The
deciding variable is measurable: the gate-active fraction - the share of steps
the safety check fires. Left bar chart: on the persistent saucepan task the gate
is active on 61.5% of steps; on the intermittent tasks, 19 to 27%. On saucepan
the gate fires *more often* than the unconditional scaler's own trigger, so
'gating' degenerates into 'always slow' - there are no safe windows to exploit.
Now the validation, and this is the part I'm proudest of: I fixed a decision rule
*in advance* - gating counts as recovering throughput only if some threshold hits
success at least 0.60 and SSM at most 0.08 - then ran the intermittent winner on
the persistent task. No row passed. The rule failed exactly where the map
predicts it should. Right plot: intermittent tasks bend down into the
safe-and-on-task corner; the persistent task just slides along the diagonal. One
picture, the whole thesis.
""")

# ============================================================================
# SLIDE 11 — Conclusion (regime map schematic + contributions)
# ============================================================================
s = slide()
top = title_bar(s, "Conclusion: a decision rule where the field had defaults")
# regime map schematic
my = top + Inches(0.1)
panel(s, Inches(0.5), my, Inches(6.0), Inches(2.55), fill=RGBColor(0xEC,0xF2,0xFB), line=BLUE)
textbox(s, Inches(0.7), my + Inches(0.14), Inches(5.6), Inches(2.3), [
    {"text": "Intermittent co-location", "size": 17, "bold": True, "color": BLUE, "space_after": 3},
    {"text": "gate active ~19\u201327% of steps", "size": 13, "italic": True, "color": GREY, "space_after": 7},
    [{"text": "\u2714 ", "size": 14, "bold": True, "color": GREEN}, {"text": "critic-gated speed-scaling", "size": 14.5, "color": INK}],
    [{"text": "\u2718 ", "size": 14, "bold": True, "color": RED}, {"text": "constrained RL: no feasible budget", "size": 14, "color": INK}],
    [{"text": "\u2718 ", "size": 14, "bold": True, "color": RED}, {"text": "binary veto: breaks chunked policy", "size": 14, "color": INK}],
])
panel(s, Inches(6.85), my, Inches(6.0), Inches(2.55), fill=RGBColor(0xFB,0xF1,0xE6), line=ORANGE)
textbox(s, Inches(7.05), my + Inches(0.14), Inches(5.6), Inches(2.3), [
    {"text": "Persistent co-location", "size": 17, "bold": True, "color": ORANGE, "space_after": 3},
    {"text": "gate active ~51\u201363% of steps", "size": 13, "italic": True, "color": GREY, "space_after": 7},
    [{"text": "\u2714 ", "size": 14, "bold": True, "color": GREEN}, {"text": "safety trained into the policy", "size": 14.5, "color": INK}],
    [{"text": "\u2718 ", "size": 14, "bold": True, "color": RED}, {"text": "reactive filters: freeze or flee", "size": 14, "color": INK}],
    [{"text": "  ", "size": 14}, {"text": "both axes only by composition (costs success)", "size": 14, "color": INK}],
])
textbox(s, Inches(0.5), my + Inches(2.65), Inches(12.33), Inches(0.5), [
    {"text": "co-location persistence  \u2014  cheap to measure on logged or simulated traffic  \u2014  predicts which mechanism wins",
     "size": 14.5, "italic": True, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
# contributions + honesty
cy = my + Inches(3.3)
textbox(s, Inches(0.5), cy, Inches(7.0), Inches(2.0), [
    {"text": "5 contributions:", "size": 15, "bold": True, "color": NAVY, "space_after": 3},
    {"text": "benchmark \u00b7 Hybrid Safety Critic \u00b7 two-axis protocol \u00b7 validated regime map \u00b7 reproducible method lessons.", "size": 13.5, "color": INK, "space_after": 8},
    {"text": "Stated honestly:", "size": 15, "bold": True, "color": RED, "space_after": 3},
    {"text": "PFL force loop wired but contact-detection blocked \u2192 claim is SSM/geometric-only; map rests on 3 tasks (persistence-dial is the named next test).", "size": 13.5, "color": INK},
])
panel(s, Inches(7.8), cy, Inches(5.03), Inches(1.95), fill=NAVY)
textbox(s, Inches(8.0), cy + Inches(0.18), Inches(4.7), Inches(1.6), [
    {"text": "Runtime filtering is not a universal safety layer.", "size": 15, "bold": True, "color": WHITE, "space_after": 6},
    {"text": "Constrained RL is not a universal alternative.", "size": 15, "bold": True, "color": WHITE, "space_after": 6},
    {"text": "The co-location pattern tells you which to use \u2014 and what it will cost.", "size": 15, "bold": True, "color": RGBColor(0x9E,0xC5,0xFF)},
])
footer(s, 11, "Conclusion")
notes(s, """
[13:15-14:45] To conclude. The headline isn't 'my method wins' - it's a decision
rule where the field had only defaults. Which safety mechanism helps is governed
by the human-robot co-location regime, not by how sophisticated the mechanism is.
Intermittent co-location: gate a speed-scaling backstop; constrained RL has no
feasible budget and a binary veto breaks the policy. Persistent co-location: only
the trained-in policy reduces exposure gracefully; reactive filters freeze or
flee. And the deciding variable - the gate-active fraction - is cheap to measure
on logged traffic before you commit. Five contributions: the benchmark, the
architecture, the two-axis protocol, the validated regime map, and the
reproducible method lessons. And I'm explicit about limits: the force side of ISO
is wired but blocked by a simulator contact-detection issue, so my compliance
claim is SSM-only; and the map rests on three tasks, which is why the named next
experiment is the persistence dial - hold the task fixed and vary only how long
the coworker lingers.
""")

# ============================================================================
# SLIDE 12 — Thanks / questions
# ============================================================================
s = slide(NAVY)
textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.0), [
    {"text": "Thank you", "size": 46, "bold": True, "color": WHITE, "space_after": 10},
    {"text": "Which safety mechanism works is decided by a property you can measure \u2014", "size": 21, "color": RGBColor(0xC9,0xD6,0xE5)},
    {"text": "how persistently the human is in the way \u2014 not by mechanism sophistication.", "size": 21, "color": RGBColor(0xC9,0xD6,0xE5)},
])
textbox(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.2), [
    {"text": "Questions welcome.", "size": 20, "bold": True, "color": WHITE, "space_after": 6},
    {"text": "github.com/ayushpatel4/safety-bigym-2", "size": 16, "color": RGBColor(0x9E,0xC5,0xFF), "name": MONO},
])
footer(s, 12, "Thank you")
notes(s, """
[14:45-15:00] Thank you - I'm happy to take questions. (Then leave the takeaway
sentence on screen: which mechanism works is decided by a measurable property -
co-location persistence - not by sophistication.) Hold the appendix slides ready
for Q&A: PID instability, velocity/joint coverage, WCSAC, proximity calibration,
perception robustness.
""")

# ============================================================================
# APPENDIX SLIDES (for Q&A) — hidden from main flow
# ============================================================================
def appendix(title, img, blurb, note, sub=None):
    s = slide()
    top = title_bar(s, title, sub, color=RGBColor(0x37,0x47,0x55))
    add_image_fit(s, img, Inches(0.5), top, Inches(8.2), Inches(5.2))
    textbox(s, Inches(8.9), top + Inches(0.1), Inches(3.95), Inches(5.0), blurb)
    footer(s, "A", "Appendix \u00b7 Q&A")
    notes(s, note)
    return s

appendix(
    "Appendix \u00b7 PID-\u03bb is seed-unstable at the feasibility boundary",
    FIG / "d0p3_3seed_lambda_regimes.png",
    [{"text": "At d=0.3 (\u2248 the task's natural cost) the 3 seeds diverge:", "size": 15, "bold": True, "color": NAVY, "space_after": 8},
     {"text": "\u03bb = 0.00 \u2192 unconstrained", "size": 14, "bullet": True, "space_after": 4},
     {"text": "\u03bb = 0.27 \u2192 graceful avoidance", "size": 14, "bullet": True, "space_after": 4},
     {"text": "\u03bb = 3.86 \u2192 windup collapse", "size": 14, "bullet": True, "space_after": 10},
     {"text": "Fix: pin \u03bb (fixed-\u03bb headline). The dual variable carries too little signal to auto-tune on the boundary.", "size": 14, "italic": True, "color": INK}],
    "Q&A: why fixed lambda. At the only feasible budget, PID has too little signal to auto-tune; three seeds give three regimes. Fixing lambda removes the ill-conditioned loop and makes the result reproducible.",
    sub="Why the headline fixes \u03bb instead of auto-tuning",
)
appendix(
    "Appendix \u00b7 fixed-\u03bb reproduces across seeds (the basin)",
    FIG / "fixlam0p1_3seed_floor075.png",
    [{"text": "All three seeds reach a sub-baseline proximity basin in mid-training.", "size": 15, "color": NAVY, "space_after": 8},
     {"text": "Operating point chosen by a SAFETY-AWARE rule: lowest deploy proximity at success \u2265 0.75.", "size": 14, "bullet": True, "space_after": 8},
     {"text": "Peak-success selection picks against the constraint \u2014 caused two false nulls before the basin was found.", "size": 14, "italic": True, "color": INK}],
    "Q&A: checkpoint selection. The avoidance is a run of consecutive checkpoints, not a noise dip; SSM co-improves inside it; it also beats the no-shaping baseline at matched success.",
    sub="Reproducibility of the \u221223% reduction",
)
appendix(
    "Appendix \u00b7 the two ISO axes, separated",
    FIG / "velocity_axis.png",
    [{"text": "Velocity (SSM-actual) axis on saucepan:", "size": 15, "bold": True, "color": NAVY, "space_after": 8},
     {"text": "graded speed-scaling is the specialist (\u221267%).", "size": 14, "bullet": True, "space_after": 4},
     {"text": "the policy helps via avoidance.", "size": 14, "bullet": True, "space_after": 4},
     {"text": "the binary veto does NOT hold the graded margin.", "size": 14, "bullet": True, "space_after": 8},
     {"text": "Each mechanism owns one axis; only the composition reaches both.", "size": 14, "italic": True, "color": INK}],
    "Q&A: division of labour on the proper ISO axis. The filter's canonical job is the velocity axis, and only a graded (not stop/go) filter delivers it.",
    sub="The filter owns velocity; the policy owns exposure",
)
appendix(
    "Appendix \u00b7 joint ISO coverage and its ceiling",
    FIG / "joint_coverage.png",
    [{"text": "Only the policy + speed-scaling composition reaches the both-safe corner.", "size": 15, "color": NAVY, "space_after": 8},
     {"text": "Cost is multiplicative: success 0.85 \u2192 0.44.", "size": 14, "bold": True, "color": RED, "bullet": True, "space_after": 8},
     {"text": "This is the persistent regime's ceiling, not a tuning defect \u2014 every separation-triggered response is active almost everywhere.", "size": 14, "italic": True, "color": INK}],
    "Q&A: why both-axis safety is expensive under persistent co-location. The reactive cost is paid on most steps and stacks with the proactive cost.",
    sub="Both axes at once = stacked cost",
)
appendix(
    "Appendix \u00b7 the coworker disruption, frame by frame",
    FIG / "fig_coworker_patrol_disruption.png",
    [{"text": "One full coworker patrol episode (illustrative reconstruction).", "size": 15, "color": NAVY, "space_after": 8},
     {"text": "Walks in \u2192 reaches at the robot \u2192 walks away \u2192 returns and reaches again.", "size": 14, "bullet": True, "space_after": 8},
     {"text": "41% of steps inside 0.3 m on this episode; min separation 0.10 m.", "size": 14, "italic": True, "color": INK}],
    "Q&A: what the disruption looks like. Three trajectory modes (patrol / walk-in / in-place) and an arm state machine with a reach gate; train and eval parameter spaces differ for OOD testing.",
    sub="The benchmark's coworker behaviour",
)
appendix(
    "Appendix \u00b7 perception robustness & generalisation",
    FIG / "separation_render_grid.png",
    [{"text": "Reduction survives noisy, lagged perception:", "size": 15, "bold": True, "color": NAVY, "space_after": 6},
     {"text": "oracle 0.236 vs noisy 0.198 \u2014 statistically indistinguishable.", "size": 14, "bullet": True, "space_after": 8},
     {"text": "Held-out gentler coworker: \u221231% (0.084 \u2192 0.058).", "size": 14, "bullet": True, "space_after": 8},
     {"text": "Worst-case tail (CVaR0.95 min-sep \u2248 0.005 m) is exogenous \u2014 unmoved by any mechanism.", "size": 14, "italic": True, "color": INK}],
    "Q&A: sim-to-real perception and OOD. Proximity threshold tau=0.3 m calibrated by contact geometry (body-surface clearance) and the empirical separation distribution.",
    sub="Noisy perception, OOD coworker, and the exogenous tail",
)

prs.save(str(OUT))
print("saved", OUT, "with", len(prs.slides._sldIdLst), "slides")
